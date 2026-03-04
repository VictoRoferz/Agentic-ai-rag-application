"""
Document Parsers — Extract text + metadata from PDF, DOCX, PPTX.
Uses PyMuPDF (fitz) for PDF, python-docx for DOCX, python-pptx for PPTX.
Docling can be enabled in settings.yaml for higher quality parsing.
"""

import re
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from loguru import logger


def parse_document(file_path: str | Path) -> dict[str, Any]:
    """
    Parse a document and return structured data.

    Returns:
        {
            "doc_id": str (UUID),
            "doc_name": str,
            "doc_type": str ("pdf"|"docx"|"pptx"),
            "pages": [
                {"page_number": int, "text": str, "section_title": str}
            ],
            "full_text": str,
            "metadata": {...}
        }
    """
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"Document not found: {path}")

    suffix = path.suffix.lower()
    logger.info(f"Parsing {path.name} ({suffix})")

    if suffix == ".pdf":
        pages = _parse_pdf(path)
    elif suffix == ".docx":
        pages = _parse_docx(path)
    elif suffix == ".pptx":
        pages = _parse_pptx(path)
    else:
        raise ValueError(f"Unsupported format: {suffix}. Supported: .pdf, .docx, .pptx")

    full_text = "\n\n".join(p["text"] for p in pages if p["text"].strip())

    result = {
        "doc_id": str(uuid.uuid4()),
        "doc_name": path.name,
        "doc_type": suffix.lstrip("."),
        "pages": pages,
        "full_text": full_text,
        "metadata": {
            "file_path": str(path.absolute()),
            "file_size_bytes": path.stat().st_size,
            "ingested_at": datetime.now(timezone.utc).isoformat(),
            "total_pages": len(pages),
            "total_chars": len(full_text),
        },
    }

    logger.info(
        f"Parsed {path.name}: {len(pages)} pages, {len(full_text)} chars"
    )
    return result


# ── PDF Parser (PyMuPDF) ──────────────────────────────────────

def _parse_pdf(path: Path) -> list[dict]:
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise ImportError("Install PyMuPDF: pip install pymupdf")

    doc = fitz.open(str(path))
    pages = []

    for page_num in range(len(doc)):
        page = doc[page_num]
        text = page.get_text("text")

        # Try to extract section title from first bold/large text
        section_title = _extract_section_title_pdf(page)

        pages.append({
            "page_number": page_num + 1,
            "text": text.strip(),
            "section_title": section_title,
        })

    doc.close()
    return pages


def _extract_section_title_pdf(page) -> str:
    """Heuristic: find the first large or bold text block as section title."""
    blocks = page.get_text("dict", flags=0)["blocks"]
    for block in blocks:
        if "lines" not in block:
            continue
        for line in block["lines"]:
            for span in line["spans"]:
                text = span["text"].strip()
                if not text or len(text) < 3:
                    continue
                # Heuristic: text larger than 12pt or bold is likely a heading
                if span["size"] > 13 or "bold" in span["font"].lower():
                    return text[:120]
    return ""


# ── DOCX Parser ───────────────────────────────────────────────

def _parse_docx(path: Path) -> list[dict]:
    try:
        from docx import Document
    except ImportError:
        raise ImportError("Install python-docx: pip install python-docx")

    doc = Document(str(path))
    pages = []
    current_page_text = []
    current_section = ""
    page_num = 1
    char_count = 0

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            current_page_text.append("")
            continue

        # Detect headings
        if para.style and para.style.name and "heading" in para.style.name.lower():
            current_section = text[:120]

        current_page_text.append(text)
        char_count += len(text)

        # Estimate page breaks (~3000 chars per page)
        if char_count >= 3000:
            pages.append({
                "page_number": page_num,
                "text": "\n".join(current_page_text).strip(),
                "section_title": current_section,
            })
            current_page_text = []
            char_count = 0
            page_num += 1

    # Remaining text
    if current_page_text:
        pages.append({
            "page_number": page_num,
            "text": "\n".join(current_page_text).strip(),
            "section_title": current_section,
        })

    return pages


# ── PPTX Parser ───────────────────────────────────────────────

def _parse_pptx(path: Path) -> list[dict]:
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("Install python-pptx: pip install python-pptx")

    prs = Presentation(str(path))
    pages = []

    for slide_num, slide in enumerate(prs.slides, 1):
        texts = []
        slide_title = ""

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    t = paragraph.text.strip()
                    if t:
                        texts.append(t)

            # Get slide title
            if shape.shape_type is not None and hasattr(shape, "text"):
                if shape == slide.shapes.title and shape.text.strip():
                    slide_title = shape.text.strip()[:120]

        # Also extract notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                texts.append(f"[Speaker Notes: {notes}]")

        pages.append({
            "page_number": slide_num,
            "text": "\n".join(texts).strip(),
            "section_title": slide_title,
        })

    return pages


# ── Batch parsing ─────────────────────────────────────────────

def parse_directory(dir_path: str | Path, supported_formats: list[str] = None) -> list[dict]:
    """Parse all supported documents in a directory."""
    if supported_formats is None:
        supported_formats = [".pdf", ".docx", ".pptx"]

    dir_path = Path(dir_path)
    results = []

    for f in sorted(dir_path.iterdir()):
        if f.suffix.lower() in supported_formats and not f.name.startswith("."):
            try:
                result = parse_document(f)
                results.append(result)
            except Exception as e:
                logger.error(f"Failed to parse {f.name}: {e}")

    logger.info(f"Parsed {len(results)} documents from {dir_path}")
    return results
